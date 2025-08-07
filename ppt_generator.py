from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from concurrent.futures import ThreadPoolExecutor
import os

def add_text_to_frame(text_frame, text):
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)

    lines = text.split("\n")
    for line in lines:
        if line.strip():
            paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = line.strip()
            run.font.size = Pt(18)

            if line.startswith("## "):
                run.font.bold = True
                run.text = line.replace("## ", "")
            elif line.startswith("### "):
                run.font.italic = True
                run.text = line.replace("### ", "")


def prepare_slide_data(index, section, text, images, inverted_image_y_positions):
    """Prepare data for each slide part, chunked and ready for slide creation."""
    max_chars_per_slide = 550
    text_parts = [text[i:i + max_chars_per_slide] for i in range(0, len(text), max_chars_per_slide)]
    slides = []

    for part_index, text_part in enumerate(text_parts):
        if index % 2 == 0:
            text_x = Inches(0.5)
            image_x = Inches(5.5)
            image_y = Inches(2.5)
        else:
            text_x = Inches(5)
            image_x = Inches(1)
            image_y = inverted_image_y_positions.get(index, Inches(2.5))

        image_filename = images[part_index % len(images)] if images else None

        slides.append({
            "section": section,
            "text_part": text_part,
            "text_x": text_x,
            "image_x": image_x,
            "image_y": image_y,
            "image_filename": image_filename,
        })

    return slides


def create_presentation(topic, sections, slides_data, images_data, inverted_image_y_positions=None):
    presentation = Presentation()
    text_top = Inches(2)
    text_width = Inches(5)
    text_height = Inches(6)

    image_width = Inches(3.5)
    image_height = Inches(3)

    if inverted_image_y_positions is None:
        inverted_image_y_positions = {}

    all_slide_content = []

    with ThreadPoolExecutor() as executor:
        futures = [
            executor.submit(
                prepare_slide_data,
                i, section, slide_text, images, inverted_image_y_positions
            )
            for i, (section, (slide_text, images)) in enumerate(zip(sections, zip(slides_data, images_data)))
        ]
        for future in futures:
            all_slide_content.extend(future.result())
    for i, slide_info in enumerate(all_slide_content):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        title_box = slide.shapes.title
        if title_box:
            section_title = slide_info["section"].upper()
            title_box.text = section_title

        content_box = slide.shapes.add_textbox(
            slide_info["text_x"], text_top, text_width, text_height
        )
        add_text_to_frame(content_box.text_frame, slide_info["text_part"])

        if slide_info["image_filename"]:
            slide.shapes.add_picture(
                slide_info["image_filename"],
                slide_info["image_x"],
                slide_info["image_y"],
                image_width,
                image_height
            )
    ppt_path = f"{topic.replace(' ', '_')}.pptx"
    presentation.save(ppt_path)
    return ppt_path
